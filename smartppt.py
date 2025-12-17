# -*- coding: utf-8 -*-
"""
PPT图片排版脚本 - 多目录独立排版版
核心逻辑：严格按文件夹顺序处理，绝不混合不同文件夹的图片。
"""

import os
import sys
import io
import configparser
from PIL import Image, ImageFile
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import math

# 修复Pillow处理大图片的问题
ImageFile.LOAD_TRUNCATED_IMAGES = True

def setup_encoding():
    """设置系统编码，解决GBK编码问题"""
    if sys.platform.startswith('win'):
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

setup_encoding()

class ConfigManager:
    """配置文件管理器"""
    
    def __init__(self, config_file='config.ini'):
        self.config_file = config_file
        self.config = configparser.ConfigParser()
        # 默认配置
        self.defaults = {
            'PATHS': {
                'image_folders': 'D:/我的图片',
                'output_ppt': '综合图片排版_A3.pptx',
                'image_area_ratio': '0.7',
                'margin': '1.5',
                'gap': '1.2'
            },
            'SETTINGS': {
                'show_filenames': 'false',
                'border_width': '1.0',
                'show_page_numbers': 'true',
                'portrait_threshold': '0.9',
                'square_min_threshold': '0.9',
                'square_max_threshold': '1.1',
                'supported_formats': '.jpg, .jpeg, .png, .bmp, .gif, .tiff, .tif, .webp, .jfif'
            }
        }
    
    def load_config(self):
        """加载配置文件，如果不存在则创建默认配置"""
        if not os.path.exists(self.config_file):
            print(f"配置文件 {self.config_file} 不存在，正在创建默认配置...")
            self.create_default_config()
            print("请修改配置文件中的路径后重新运行程序")
            return False
        
        try:
            # 明确指定UTF-8编码读取，避免中文路径问题
            self.config.read(self.config_file, encoding='utf-8')
            print(f"成功加载配置文件: {self.config_file}")
            return True
        except Exception as e:
            print(f"读取配置文件失败: {e}")
            return False
    
    def create_default_config(self):
        """创建默认配置文件"""
        try:
            self.config.read_dict(self.defaults)
            with open(self.config_file, 'w', encoding='utf-8') as f:
                self.config.write(f)
            print(f"已创建默认配置文件: {self.config_file}")
        except Exception as e:
            print(f"创建配置文件失败: {e}")
    
    def get_image_folders(self):
        """获取图片文件夹列表"""
        try:
            folders_str = self.config['PATHS']['image_folders']
            folders = []
            for folder in folders_str.split(','):
                folder = folder.strip()
                if folder:
                    # 处理路径中的环境变量和用户目录
                    folder = os.path.expandvars(folder)
                    folder = os.path.expanduser(folder)
                    folder = os.path.abspath(folder)
                    folders.append(folder)
            return folders
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError):
            # 如果配置项不存在，返回默认值
            default_folders = [self.defaults['PATHS']['image_folders']]
            return [os.path.abspath(folder) for folder in default_folders]
    
    def get_output_path(self):
        """获取输出文件路径"""
        try:
            path = self.config['PATHS']['output_ppt']
            path = os.path.expandvars(path)
            path = os.path.expanduser(path)
            path = os.path.abspath(path)
            return path
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError):
            return os.path.abspath(self.defaults['PATHS']['output_ppt'])
    
    def get_float(self, section, key):
        """获取浮点数配置"""
        try:
            return float(self.config[section][key])
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError, ValueError):
            return float(self.defaults[section][key])
    
    def get_bool(self, section, key):
        """获取布尔值配置"""
        try:
            value = self.config[section][key].lower()
            return value in ['true', 'yes', '1', 'on']
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError):
            value = self.defaults[section][key].lower()
            return value in ['true', 'yes', '1', 'on']
    
    def get_list(self, section, key):
        """获取列表配置"""
        try:
            value = self.config[section][key]
            return [ext.strip() for ext in value.split(',')]
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError):
            value = self.defaults[section][key]
            return [ext.strip() for ext in value.split(',')]
    
    def get_string(self, section, key):
        """获取字符串配置"""
        try:
            return self.config[section][key]
        except (KeyError, configparser.NoSectionError, configparser.NoOptionError):
            return self.defaults[section][key]

class SmartPPTGenerator:
    def __init__(self, config_manager):
        self.config = config_manager
        self.output_path = self.config.get_output_path()
        self.image_area_ratio = self.config.get_float('PATHS', 'image_area_ratio')
        self.page_width = Cm(42)
        self.page_height = Cm(29.7)
        self.margin = Cm(self.config.get_float('PATHS', 'margin'))
        self.gap = Cm(self.config.get_float('PATHS', 'gap'))
        self.show_filenames = self.config.get_bool('SETTINGS', 'show_filenames')
        self.border_width = Pt(self.config.get_float('SETTINGS', 'border_width'))
        self.show_page_numbers = self.config.get_bool('SETTINGS', 'show_page_numbers')
        self.supported_formats = tuple(self.config.get_list('SETTINGS', 'supported_formats'))
        self.portrait_threshold = self.config.get_float('SETTINGS', 'portrait_threshold')
        self.square_min = self.config.get_float('SETTINGS', 'square_min_threshold')
        self.square_max = self.config.get_float('SETTINGS', 'square_max_threshold')

    def safe_print(self, *args, **kwargs):
        """安全的打印函数"""
        try:
            print(*args, **kwargs)
        except UnicodeEncodeError:
            try:
                text = ' '.join(str(arg) for arg in args)
                encoded = text.encode('utf-8', errors='replace').decode('utf-8')
                print(encoded, **kwargs)
            except:
                print("[编码错误]", **kwargs)

    def robust_read_image(self, filepath):
        """健壮的图片读取方法"""
        filename = os.path.basename(filepath)
        try:
            with Image.open(filepath) as img:
                img.load()
                return {
                    'path': filepath,
                    'filename': filename,
                    'width': img.size[0],
                    'height': img.size[1],
                    'ratio': img.size[0] / img.size[1],
                    'success': True
                }
        except Exception:
            return None

    def classify_images_in_folder(self, folder_path):
        """分类单个文件夹中的图片"""
        portraits = []
        squares = []
        failed = 0

        for filename in os.listdir(folder_path):
            if any(filename.lower().endswith(ext) for ext in self.supported_formats):
                filepath = os.path.join(folder_path, filename)
                img_info = self.robust_read_image(filepath)
                if img_info:
                    ratio = img_info['ratio']
                    if ratio <= self.portrait_threshold:
                        portraits.append(img_info)
                    elif self.square_min < ratio < self.square_max:
                        squares.append(img_info)
                    else:  # 横版归为方形
                        squares.append(img_info)
                else:
                    failed += 1
        return portraits, squares, failed

    def create_mixed_slide(self, slide, square_img, portrait_img):
        """创建混合页面（1方+1竖）"""
        available_width = self.page_width - (2 * self.margin)
        available_height = self.page_height - (2 * self.margin)
        unified_height = available_height * self.image_area_ratio

        square_width = unified_height * square_img['ratio']
        portrait_width = unified_height * portrait_img['ratio']
        total_width = square_width + self.gap + portrait_width

        if total_width > available_width:
            scale = available_width / total_width
            square_width *= scale
            portrait_width *= scale
            unified_height *= scale
            total_width = square_width + self.gap + portrait_width

        start_x = self.margin + (available_width - total_width) / 2
        start_y = self.margin + (available_height - unified_height) / 2

        # 添加图片
        self.add_image_to_slide(slide, square_img, start_x, start_y, square_width, unified_height)
        self.add_image_to_slide(slide, portrait_img, start_x + square_width + self.gap, start_y, portrait_width, unified_height)

    def create_portrait_slide(self, slide, portrait_imgs):
        """创建纯竖版页面（3张竖版）"""
        available_width = self.page_width - (2 * self.margin)
        available_height = self.page_height - (2 * self.margin)
        unified_height = available_height * self.image_area_ratio

        image_widths = [unified_height * img['ratio'] for img in portrait_imgs]
        total_gap = (len(portrait_imgs) - 1) * self.gap
        total_width = sum(image_widths) + total_gap

        if total_width > available_width:
            scale = available_width / total_width
            unified_height *= scale
            image_widths = [unified_height * img['ratio'] for img in portrait_imgs]
            total_width = sum(image_widths) + total_gap

        start_x = self.margin + (available_width - total_width) / 2
        start_y = self.margin + (available_height - unified_height) / 2

        current_x = start_x
        for i, img in enumerate(portrait_imgs):
            self.add_image_to_slide(slide, img, current_x, start_y, image_widths[i], unified_height)
            current_x += image_widths[i] + self.gap

    def add_image_to_slide(self, slide, img_info, left, top, width, height):
        """添加图片到幻灯片"""
        try:
            orig_ratio = img_info['ratio']
            if orig_ratio >= 1:
                new_width = min(width, height * orig_ratio)
                new_height = new_width / orig_ratio
            else:
                new_height = min(height, width / orig_ratio)
                new_width = new_height * orig_ratio

            left_c = left + (width - new_width) / 2
            top_c = top + (height - new_height) / 2

            pic = slide.shapes.add_picture(img_info['path'], left_c, top_c, new_width, new_height)
            if self.border_width.pt > 0:
                pic.line.width = self.border_width
                pic.line.color.rgb = RGBColor(180, 180, 180)
            return True
        except Exception as e:
            self.safe_print(f"  添加图片失败 {img_info['filename']}: {e}")
            return False

    def add_page_number(self, slide, page_num):
        """添加页码"""
        if not self.show_page_numbers:
            return
        try:
            txBox = slide.shapes.add_textbox(self.page_width - Cm(2.5), self.page_height - Cm(1.0), Cm(2.0), Cm(0.6))
            tf = txBox.text_frame
            tf.text = f"{page_num}"
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        except:
            pass

    def generate_ppt(self):
        """生成PPT - 核心修改：按文件夹独立处理"""
        image_folders = self.config.get_image_folders()
        if not image_folders:
            self.safe_print("错误: 未在配置中找到有效的图片文件夹路径。")
            return

        prs = Presentation()
        prs.slide_width = self.page_width
        prs.slide_height = self.page_height
        blank_layout = prs.slide_layouts[6]

        total_slide_count = 0
        folder_summary = []

        self.safe_print("=" * 70)
        self.safe_print("开始按文件夹顺序独立排版...")
        self.safe_print("=" * 70)

        for folder_index, folder_path in enumerate(image_folders):
            folder_name = os.path.basename(os.path.normpath(folder_path))
            self.safe_print(f"\n>>> 正在处理第 {folder_index + 1} 个文件夹: {folder_name}")

            if not os.path.exists(folder_path):
                self.safe_print(f"  警告: 文件夹不存在，跳过。")
                folder_summary.append((folder_name, "文件夹不存在", 0))
                continue

            # 1. 分类当前文件夹的图片
            portraits, squares, failed = self.classify_images_in_folder(folder_path)
            total_imgs = len(portraits) + len(squares)

            self.safe_print(f"  扫描结果: 共 {total_imgs} 张可用图片 ({len(portraits)} 竖, {len(squares)} 方/横), {failed} 张读取失败。")

            if total_imgs == 0:
                self.safe_print(f"  提示: 文件夹内无可用图片，跳过。")
                folder_summary.append((folder_name, "无图片", 0))
                continue

            # 2. 处理当前文件夹的图片
            folder_slide_count = 0
            square_idx, portrait_idx = 0, 0

            # 2.1 先排混合页 (1方 + 1竖)
            while square_idx < len(squares) and portrait_idx < len(portraits):
                slide = prs.slides.add_slide(blank_layout)
                total_slide_count += 1
                folder_slide_count += 1
                self.create_mixed_slide(slide, squares[square_idx], portraits[portrait_idx])
                self.add_page_number(slide, total_slide_count)
                square_idx += 1
                portrait_idx += 1

            # 2.2 再排剩余的竖版图片 (3张一页)
            remaining_portraits = len(portraits) - portrait_idx
            if remaining_portraits > 0:
                portrait_groups = math.ceil(remaining_portraits / 3)
                for group in range(portrait_groups):
                    slide = prs.slides.add_slide(blank_layout)
                    total_slide_count += 1
                    folder_slide_count += 1
                    start = portrait_idx + (group * 3)
                    end = min(start + 3, len(portraits))
                    self.create_portrait_slide(slide, portraits[start:end])
                    self.add_page_number(slide, total_slide_count)

            # 记录当前文件夹的统计
            unused_squares = len(squares) - square_idx
            status_note = f"已用 {min(len(squares), len(portraits))} 方" + (f", {unused_squares} 方未匹配" if unused_squares > 0 else "")
            folder_summary.append((folder_name, status_note, folder_slide_count))
            self.safe_print(f"  完成: 生成了 {folder_slide_count} 页。{status_note}")

        # 3. 保存并输出总结果
        self.safe_print("\n" + "=" * 70)
        self.safe_print("所有文件夹处理完成！")
        self.safe_print("=" * 70)
        self.safe_print("各文件夹排版摘要:")
        for name, note, pages in folder_summary:
            self.safe_print(f"  · {name}: {pages} 页 ({note})")
        self.safe_print(f"PPT总页数: {total_slide_count}")

        try:
            prs.save(self.output_path)
            file_size = os.path.getsize(self.output_path)
            self.safe_print(f"文件已保存: {os.path.abspath(self.output_path)} ({file_size/1024/1024:.1f} MB)")
            if sys.platform.startswith('win') and os.path.exists(self.output_path):
                try:
                    os.startfile(self.output_path)
                    self.safe_print("已尝试自动打开PPT文件。")
                except:
                    pass
        except Exception as e:
            self.safe_print(f"\n保存PPT时出错: {e}")

def main():
    print("PPT图片排版脚本 - 文件夹独立排版版")
    print("=" * 70)
    config_manager = ConfigManager('config.ini')
    if not config_manager.load_config():
        print("使用默认配置，请修改config.ini文件。")
    print(f"将按顺序处理 {len(config_manager.get_image_folders())} 个文件夹。")
    print("=" * 70)
    input("按回车键开始...")
    generator = SmartPPTGenerator(config_manager)
    generator.generate_ppt()
    input("\n按回车键退出...")

if __name__ == "__main__":
    main()